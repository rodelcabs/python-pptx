import settings
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.dml import MSO_THEME_COLOR_INDEX

def getSlideTextList (contentFileName):
    newContent = []

    with open(settings.srcConntents+contentFileName) as f:
        lines = f.readlines()
        currentTxt = ""

        for x, txt in enumerate(lines):
            
            if(txt == "\n"):
                currentTxt = ""
                continue
            else:
                currentTxt += txt

                if((x+1) == len(lines) or ((x+1) < len(lines) and lines[x+1] == "\n")):
                   newContent.append(currentTxt) 
           
    return [x.replace("\n", " ") for x in newContent]


def createPPT(bgName, title, fileName, contentList, existing=False):
    backgroundPic = settings.srcConntents+bgName

    # creates new ppt presentation
    if existing:
        prs = Presentation(settings.existingPPTPath+fileName)
    else:
        prs = Presentation()

    blankSlide = prs.slide_layouts[6]

    # initial slide for title
    initialSlide = prs.slides[0] if existing else prs.slides.add_slide(blankSlide)
    initialSlide.shapes.add_picture(
            backgroundPic, 
            Inches(0), 
            Inches(0), 
            height=Inches(settings.bgPicHeight), 
            width=Inches(settings.bgPicWidth)
        )
    titleTextBox = initialSlide.shapes.add_textbox(
            Inches(settings.titleTextBox["top"]),
            Inches(settings.titleTextBox["left"]),
            Inches(settings.titleTextBox["width"]),
            Inches(settings.titleTextBox["height"]),
        )
    tf = titleTextBox.text_frame
    titleTXT = tf.add_paragraph()
    titleTXT.text = title
    titleTXT.font.color.theme_color = MSO_THEME_COLOR_INDEX["ACCENT_1"]
    titleTXT.font.size = Pt(settings.titleTextBox["fontSize"])
    titleTXT.font.name = "Arial Black"

    prs.save(settings.destination+fileName)


